import os
import glob
import pptx
import csv
import pandas as pd
import filetype
# import base64
import moviepy.editor as mp
import speech_recognition as sr
from pydub import AudioSegment
import extract_msg
import openpyxl
import pyxlsb
import docx
import xlrd
from langchain.text_splitter import RecursiveCharacterTextSplitter
from brain import embeddings
import fitz
from langchain_community.vectorstores.chroma import Chroma

CHROMA_PATH = "Database"


# # Clear out the database first
# if os.path.exists(CHROMA_PATH):
#     shutil.rmtree(CHROMA_PATH)


def generate_tokens(s, file_path):
    l=1000
    n = len(s)
    if n <= 1000:
        l = 200
    elif n <=2500:
        l = 500
    elif n<=5000:
        l = 1000
    elif n<=7500:
        l = 1500
    else:
        l = 2000

    text_splitter = RecursiveCharacterTextSplitter(        
        separators=["\n\n", "\n", ".", " "],
        chunk_size = l,
        chunk_overlap  = l//10,
        length_function = len,
        )
    splits = text_splitter.split_text(s)
    documents = text_splitter.create_documents(splits)
    for doc in documents:
        doc.metadata['file_path'] = file_path
    Chroma.from_documents(documents, embeddings, persist_directory = CHROMA_PATH)


def extract_text_from_file(file_path, img_counter):
    file_extension = os.path.splitext(file_path)[1].lower()
    try:
            
        if file_extension == '.docx':
            doc = docx.Document(file_path)
            raw_text = ''
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_counter += 1
                    image_blob = rel.target_part.blob
                    image_format = rel.target_part.content_type.split('/')[-1]
                    with open(os.path.join("output_images", f"image_{img_counter}.png"), "wb") as f:
                        f.write(image_blob)
            
            for paragraph in doc.paragraphs:
                raw_text += paragraph.text + "\n"
            
            if raw_text.strip():
                generate_tokens(raw_text, file_path=file_path)
            return raw_text


        elif file_extension == '.pptx':

            raw_text = ''
            prs = pptx.Presentation(file_path)
            
            for slide_number, slide in enumerate(prs.slides): 
                for shape in slide.shapes: 
                    if hasattr(shape, "text"): 
                        raw_text += (shape.text + "\n")
                    elif hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        with open(os.path.join("output_images", f"image_{img_counter}.png"), 'wb') as f:
                            f.write(image_bytes)
                        raw_text += f"\n</IMAGE_{img_counter}>\n"
                        img_counter += 1
            
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
        
        # RAW_TEXT IS IN UNICODE SO DON'T PRINT OR TRY TO ENCODE
        # IF WRITING IN TEXT FILE THEN USE ENCODING = "UTF-8"
        elif file_extension == '.pdf':

            doc = fitz.open(file_path)
            raw_text = ""
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text("text")
                raw_text += (text + "\n")
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    with open(os.path.join("output_images", f"image_{img_counter}.png"), 'wb') as img_file:
                        img_file.write(image_bytes)
                    raw_text += f"\n</IMAGE_{img_counter}> \n"
                    img_counter += 1
            
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
        
        elif file_extension == '.txt':
            raw_text = ''
            with open(file_path, 'r', encoding='utf-8') as f:
                raw_text += f.read()
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
        
        elif file_extension == '.xlsx' or file_extension == '.xlsm':
            raw_text = ""
            xlsx_file = pd.ExcelFile(file_path)
            
            for sheet_name in xlsx_file.sheet_names:
                raw_text += f"Sheet: {sheet_name}:\n"
                sheet_data = pd.read_excel(xlsx_file, sheet_name=sheet_name, engine='openpyxl', header=None)
                # Iterate through each table in the sheet (assuming each table is separated by at least one NaN row)
                for _, table in sheet_data.groupby(sheet_data.isnull().all(1)):
                    if not table.empty:
                        # Convert the table to text (assuming it's already in DataFrame format)
                        table_text = table.to_string(index=False, header=False)
                        raw_text +=table_text + " "
                    raw_text +="\n"
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text

        elif file_extension == '.xlsb':
            raw_text = ''
            wb = pyxlsb.open_workbook(file_path)

            for sheet_name in wb.sheets:
                raw_text += f"Sheet: {sheet_name}:\n"
                with wb.get_sheet(sheet_name) as sheet:
                    for row in sheet.rows():
                        for item in row:
                            if item.v is not None:
                                raw_text += str(item.v) + " "
                    raw_text += '\n'
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
            



        elif file_extension == '.xls':
            raw_text = ""
            wb = xlrd.open_workbook(file_path)
            for sheet_name in wb.sheet_names():
                raw_text += f"Sheet: {sheet_name}:\n"
                sheet = wb.sheet_by_name(sheet_name)
                
                # Iterate through each row and column to extract text
                for row_idx in range(sheet.nrows):
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        raw_text += cell_value + " "
                    raw_text += '\n'
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
        
            
        
        elif file_extension == '.csv':
            raw_text = ""
            with open(file_path, 'r', newline='') as csvfile:
                reader = csv.DictReader(csvfile)
                for row in reader:
                    cell_texts = []
                    for key, value in row.items():
                        if value:
                            cell_texts.append(f"{key}: {value}")
                    cell_text = ", ".join(cell_texts)
                    raw_text += cell_text + "\n"
            if raw_text:
                generate_tokens(raw_text, file_path=file_path)
            return raw_text
        
        elif file_extension == '.msg':
            raw_text=''
            msg = extract_msg.Message(file_path)
            raw_text = msg.body
            if raw_text:
                generate_tokens(raw_text,file_path=file_path)
            return raw_text
        
        #Check for json string to see if it can be further improved.
        elif filetype.is_image(file_path):
            with open(os.path.join("output_images", f"image_{img_counter}.png"), 'wb') as img_file:
                img_file.write(file_path)
            return f"\n</IMAGE_{img_counter}> \n"
            # image_url = f"data:image/png;base64,{encode_image(file_path)}"
            
            # json_string = ai_image_response(image_url=image_url).choices[0].message.content
            # json_string = json_string.replace("```json\n","").replace("\n```","")

            # return json_string
            
        # elif filetype.is_video(file_path):
        #     audio_path="output_audio.wav"
        #     video = mp.VideoFileClip(file_path)
        #     audio = video.audio
        #     audio.write_audiofile(audio_path)
        #     recognizer = sr.Recognizer()
        #     audio_file = sr.AudioFile(audio_path, language='en-US', endpoint='https://www.google.com/speech-api/v2/recognize')
        #     with audio_file as source:
        #         recognizer.adjust_for_ambient_noise(source)
        #         audio = recognizer.record(source)
        #     raw_text = ''
        #     try:
        #         raw_text = recognizer.recognize_google(audio)
        #         if raw_text:
        #                 generate_tokens(raw_text, file_path=file_path)
        #         return raw_text
        #     except sr.UnknownValueError:
        #         print("Could not understand audio")
        #     except sr.RequestError:
        #         print("Could not request results")


        elif filetype.is_audio(file_path):
            recognizer = sr.Recognizer()
            with sr.AudioFile(file_path) as source:
                audio_data = recognizer.record(source)
            
            try:
                raw_text = recognizer.recognize_google(audio_data)
                if raw_text:
                    generate_tokens(raw_text, file_path=file_path)
                return raw_text
        
            except sr.UnknownValueError:
                print("Could not understand audio")
            except sr.RequestError as e:
                print(f"could not request results; {e}")

        # # IMPORT WHISPER FILE BY EXTENSION OR DOWNLOAD
        # # IF THERE IS A TENSOR ERROR THEN CHANGE TRANSCRIBE FUNCTION AS SUCH
        # # model.transcibe(file_path, fp16=False)
        # elif filetype.is_audio(file_path):
        #     model = whisper.load_model('base')
        #     result = model.transcribe(file_path)
        #     return result['text']
        # Add more conditions for other file types if needed
        else:
            print("Skipped")
            return ''
    except Exception as e:
        print("ERROR:>>>>>>>>>")
        print(file_path)
        print(e)
        return ''
def extract_text_from_folder(folder_path):
    all_text = ''
    img_counter =0
    for file_path in glob.glob(os.path.join(folder_path, '**'), recursive=True):
        if os.path.isfile(file_path):
            text = extract_text_from_file(file_path, img_counter)
            print(file_path)
            if text:
                all_text += text
    return all_text



path = r''
exception = []
raw_text = extract_text_from_folder(path)
                
with open(r'Output.txt', 'a',  encoding='utf-8') as f:
    f.write(raw_text)
